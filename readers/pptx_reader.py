from io import BytesIO

from pptx import Presentation

from .base import BaseReader


class PptxReader(BaseReader):

    extensions = ["pptx"]

    def read(self, filename, contents):

        prs = Presentation(BytesIO(contents))

        slides = []

        for i, slide in enumerate(prs.slides, 1):

            slides.append(f"===== Slide {i} =====")

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slides.append(shape.text)

        return "\n".join(slides)